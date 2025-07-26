import os
from typing import List
from pypdf import PdfReader
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def extract_text_from_pdf(file_path: str) -> str:
    """
    Extracts text from a PDF file.

    Args:
        file_path (str): The path to the PDF file.

    Returns:
        str: The extracted text content.
    """
    text = ""
    try:
        reader = PdfReader(file_path)
        for page in reader.pages:
            # Extract text from each page, handling potential None returns
            text += page.extract_text() or ""
    except Exception as e:
        print(f"Error extracting text from PDF {file_path}: {e}")
    return text

def extract_text_from_pptx(file_path: str) -> str:
    """
    Extracts text from a PPTX file.

    Args:
        file_path (str): The path to the PPTX file.

    Returns:
        str: The extracted text content.
    """
    text = ""
    try:
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                # Check if the shape has text directly
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
                # If it's a group, iterate through its sub-shapes
                elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    for s in shape.shapes:
                        if hasattr(s, "text"):
                            text += s.text + "\n"
    except Exception as e:
        print(f"Error extracting text from PPTX {file_path}: {e}")
    return text

def chunk_text(text: str, chunk_size: int = 500, overlap: int = 50) -> List[str]:
    """
    Chunks text into smaller pieces with optional overlap.
    This helps in managing context for the LLM and vector database.

    Args:
        text (str): The input text to be chunked.
        chunk_size (int): The maximum number of words per chunk.
        overlap (int): The number of words to overlap between consecutive chunks.

    Returns:
        List[str]: A list of text chunks.
    """
    chunks = []
    if not text:
        return chunks

    words = text.split() # Simple word-based splitting
    current_chunk = []
    for word in words:
        current_chunk.append(word)
        # If the current chunk reaches the desired size, add it and prepare for overlap
        if len(current_chunk) >= chunk_size:
            chunks.append(" ".join(current_chunk))
            # Keep the last 'overlap' words for the next chunk
            current_chunk = current_chunk[chunk_size - overlap:]
    # Add any remaining words as the last chunk
    if current_chunk:
        chunks.append(" ".join(current_chunk))
    return chunks

